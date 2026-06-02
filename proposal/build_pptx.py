"""Assemble an editable, Canva-importable .pptx from the HTML-derived assets.

Reads layout.json (text blocks with positions/styles) + backgrounds/bg-N.png
(text-less page images) produced by render.mjs, and writes
SNEC_EyeBot_Proposal.pptx — A4 portrait, with every text block as a real,
editable PowerPoint text box over a pixel-perfect background.
"""
import json
import os

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
EMU_PER_PX = 914400 / 96  # 96 CSS px per inch

ALIGN = {
    "left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT, "end": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def px(v):
    return Emu(int(round(v * EMU_PER_PX)))


def split_paragraphs(runs):
    """Turn a flat run list (with '\n' markers / embedded newlines) into paragraphs."""
    paras = [[]]
    for r in runs:
        parts = r["text"].split("\n")
        for j, part in enumerate(parts):
            if j > 0:
                paras.append([])
            if part != "":
                paras[-1].append({**r, "text": part})
    return paras


def main():
    layout = json.load(open(os.path.join(HERE, "layout.json"), encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Emu(int(794 * EMU_PER_PX))
    prs.slide_height = Emu(int(1123 * EMU_PER_PX))
    blank = prs.slide_layouts[6]

    for i, els in enumerate(layout):
        slide = prs.slides.add_slide(blank)
        bg = os.path.join(HERE, "backgrounds", f"bg-{i + 1}.png")
        if os.path.exists(bg):
            slide.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)

        for el in els:
            tb = slide.shapes.add_textbox(px(el["x"]), px(el["y"]), px(el["w"]), px(el["h"]))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            try:
                tf.vertical_anchor = MSO_ANCHOR.TOP
            except Exception:
                pass

            al = ALIGN.get(el.get("align", "left"), PP_ALIGN.LEFT)
            first = True
            for prun in split_paragraphs(el["runs"]):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = al
                try:
                    p.line_spacing = Pt(el["linePt"])
                except Exception:
                    pass
                for r in prun:
                    run = p.add_run()
                    run.text = r["text"]
                    f = run.font
                    f.name = el["font"]
                    f.size = Pt(max(el["sizePt"], 1))
                    f.bold = bool(r.get("bold"))
                    f.italic = bool(r.get("italic"))
                    c = r.get("color", [0, 0, 0])
                    f.color.rgb = RGBColor(c[0], c[1], c[2])

    out = os.path.join(HERE, "SNEC_EyeBot_Proposal.pptx")
    prs.save(out)
    print(f"pptx written: {out}  ({len(layout)} slides)")


if __name__ == "__main__":
    main()
